from pathlib import Path
import sys

from pptx import Presentation


def main() -> int:
    if len(sys.argv) != 2:
        print("Usage: create_base_pptx.py <output_path>", file=sys.stderr)
        return 1

    output_path = Path(sys.argv[1])
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()

    title_layout = presentation.slide_layouts[0]
    title_slide = presentation.slides.add_slide(title_layout)
    title_slide.shapes.title.text = "Working Deck"
    title_slide.placeholders[1].text = "Skill runtime starter presentation"

    content_layout = presentation.slide_layouts[1]
    content_slide = presentation.slides.add_slide(content_layout)
    content_slide.shapes.title.text = "Key Points"
    content_slide.placeholders[1].text = "Replace this content with the requested narrative."

    presentation.save(str(output_path))
    print(output_path)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
